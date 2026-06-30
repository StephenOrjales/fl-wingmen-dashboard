"""
District Report PPT Generator for FL Wingmen Dashboard.
Generates a branded PowerPoint deck per district for weekly DM meetings.
Color-coded cells match dashboard thresholds. Scorecard is the concluding summary.
"""
import io
import re
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DATA_DIR = Path(__file__).parent / "data"

# ── Brand colors ──
DARK_GREEN = RGBColor(0x1A, 0x3C, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
BLACK = RGBColor(0x1F, 0x29, 0x37)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _add_title_bar(slide, title_text, subtitle_text=""):
    """Add a dark green title bar at the top of the slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_GREEN
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xBB)
        p2.alignment = PP_ALIGN.LEFT


def _add_table(slide, df, left, top, width, height, header_color=DARK_GREEN,
               cell_colors=None, bold_cells=None, col_widths=None):
    """Add a styled table to the slide.

    cell_colors: dict mapping (row_idx, col_idx) -> RGBColor for font color
    bold_cells:  dict mapping (row_idx, col_idx) -> True for bold font
    col_widths:  list of relative widths (e.g. [2, 1, 1, 1]) — auto-normalised
    """
    rows, cols = len(df) + 1, len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths and len(col_widths) == cols:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(width * cw / total)
    else:
        col_w = int(width / cols)
        for i in range(cols):
            table.columns[i].width = col_w

    # Header row
    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    cell_colors = cell_colors or {}
    bold_cells = bold_cells or {}
    for i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if pd.notna(val) else "—"
            bg = LIGHT_GRAY if i % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            font_color = cell_colors.get((i, j), BLACK)
            is_bold = bold_cells.get((i, j), False)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = font_color
                p.font.bold = is_bold
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def _add_callouts(slide, callouts, top):
    """Add a 'Key Callouts' box."""
    left = Inches(0.5)
    width = SLIDE_WIDTH - Inches(1)
    height = Inches(0.3 + 0.25 * len(callouts))

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)  # light yellow
    box.line.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "Key Callouts"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)

    for c in callouts:
        p2 = tf.add_paragraph()
        p2.text = f"• {c}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = BLACK


def _add_section_label(slide, text, left, top, width):
    """Add a section sub-header label on a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.LEFT


def _fmt_sos(minutes):
    if pd.isna(minutes):
        return "—"
    m = int(minutes)
    s = int(round((minutes - m) * 60))
    return f"{m}:{s:02d}"


def _get_district_stores(district, store_to_district):
    """Get list of store numbers for a district."""
    return [snum for snum, d in store_to_district.items() if d == district]


def _add_scorecard_slide(prs, blank_layout, district, d_stores, sc_qtr, sc_qtr_ps,
                         check_labels, kds, labor, smg, qi, fl, cogs):
    """Build one Scorecard slide for the given quarter. Returns the slide (or None if no stores)."""
    sc_data = {}
    for snum in d_stores:
        chk = {}
        if kds is not None:
            kq = kds[(kds["_p"].isin(sc_qtr_ps)) & (kds["Store No"] == snum)]
            if not kq.empty:
                a_sos, a_ad = kq["SOS"].mean(), kq["Adoption %"].mean()
                a_ma, a_pb = kq["Make Ahead %"].mean(), kq["Pre-Bump %"].mean()
                chk["KDS: SOS < 10"] = a_sos < 10 if pd.notna(a_sos) else None
                chk["KDS: Adopt ≥ 85%"] = a_ad >= 85 if pd.notna(a_ad) else None
                chk["KDS: MakeAhd ≤ 10%"] = a_ma <= 10 if pd.notna(a_ma) else None
                chk["KDS: PreBump ≤ 0.5%"] = a_pb <= 0.5 if pd.notna(a_pb) else None
        if labor is not None:
            lab = labor[(labor["period"].isin(sc_qtr_ps)) & (labor["store_num"] == snum)]
            if not lab.empty:
                tg, th = lab["guide_hours"].sum(), lab["actual_crew_hours"].sum()
                if tg > 0:
                    chk["Labor: Var ≤ 1.5%"] = ((th - tg) / tg * 100) <= 1.5
        if smg is not None:
            sk = smg[smg["Store No"] == snum]
            if not sk.empty:
                r = sk.iloc[0]
                chk["SMG: Dissat ≤ 3%"] = r["Dissatisfaction %"] <= 3 if pd.notna(r.get("Dissatisfaction %")) else None
                chk["SMG: Inacc ≤ 3%"] = r["Inaccurate Order %"] <= 3 if pd.notna(r.get("Inaccurate Order %")) else None
        if qi is not None:
            sk = qi[qi["Store No"] == snum]
            if not sk.empty:
                val = sk.iloc[0].get("QSC Stars")
                chk["QSC: 5 Stars"] = val >= 5 if pd.notna(val) else None
        if fl is not None:
            sk = fl[fl["Store No"] == snum]
            if not sk.empty:
                val = sk.iloc[0].get("Avg_Completion_Rate")
                chk["FlavorLab: ≥ 95%"] = val >= 95 if pd.notna(val) else None
        if cogs is not None:
            cq = cogs[(cogs["_p"].isin(sc_qtr_ps)) & (cogs["Store No"] == snum)]
            if not cq.empty:
                acv = cq["COGS Variance %"].mean()
                chk["COGS: Var ≤ 1%"] = abs(acv) <= 1 if pd.notna(acv) else None
        sc_data[snum] = chk

    if not sc_data:
        return None
    slide = prs.slides.add_slide(blank_layout)
    _add_title_bar(slide, f"Scorecard — Q{sc_qtr} QTD Adherence Summary",
                   f"{district} · Concluding Performance Overview")

    sc_rows = []
    for snum in sorted(sc_data.keys(), key=int):
        row_data = {"Store #": snum}
        passed = total = 0
        for cl in check_labels:
            result = sc_data[snum].get(cl)
            if result is not None and bool(result):
                row_data[cl] = "✅"; passed += 1; total += 1
            elif result is not None and not bool(result):
                row_data[cl] = "❌"; total += 1
            else:
                row_data[cl] = "—"
        adh = (passed / total * 100) if total > 0 else 0
        row_data["Adherence"] = f"{adh:.0f}%"
        row_data["_adh_raw"] = adh
        row_data["_passed"] = passed
        row_data["_total"] = total
        sc_rows.append(row_data)

    sc_df = pd.DataFrame(sc_rows)
    display_cols = ["Store #"] + check_labels + ["Adherence"]
    sc_display = sc_df[display_cols].copy()

    colors, bolds = {}, {}
    col_names = list(sc_display.columns)
    for i, (_, row) in enumerate(sc_df.iterrows()):
        for cl in check_labels:
            j = col_names.index(cl)
            result = sc_data[row["Store #"]].get(cl)
            if result is not None and bool(result):
                colors[(i, j)] = GREEN
            elif result is not None and not bool(result):
                colors[(i, j)] = RED; bolds[(i, j)] = True
        adh = row["_adh_raw"]
        j_adh = col_names.index("Adherence")
        if adh >= 90:
            colors[(i, j_adh)] = GREEN; bolds[(i, j_adh)] = True
        elif adh < 50:
            colors[(i, j_adh)] = RED; bolds[(i, j_adh)] = True

    cw = [1.5] + [1] * len(check_labels) + [1.3]
    table_h = 0.55 + 0.32 * len(sc_display)
    table_top = 1.4
    _add_table(slide, sc_display, Inches(0.3), Inches(table_top), SLIDE_WIDTH - Inches(0.6),
               Inches(table_h), cell_colors=colors, bold_cells=bolds, col_widths=cw)

    callouts = []
    adherences = [r["_adh_raw"] for r in sc_rows]
    avg_adh = sum(adherences) / len(adherences) if adherences else 0
    callouts.append(f"District avg adherence: {avg_adh:.0f}%")
    for cl in check_labels:
        failed = sum(1 for r in sc_rows if sc_data[r["Store #"]].get(cl) is not None
                     and not bool(sc_data[r["Store #"]].get(cl)))
        twd = sum(1 for r in sc_rows if sc_data[r["Store #"]].get(cl) is not None)
        if twd > 0 and failed / twd > 0.5:
            callouts.append(f"{cl}: {failed}/{twd} stores failing — needs district focus")
    lowest = min(sc_rows, key=lambda r: r["_adh_raw"])
    callouts.append(f"Lowest: Store {lowest['Store #']} at {lowest['_adh_raw']:.0f}% ({lowest['_passed']}/{lowest['_total']})")
    highest = max(sc_rows, key=lambda r: r["_adh_raw"])
    callouts.append(f"Highest: Store {highest['Store #']} at {highest['_adh_raw']:.0f}% ({highest['_passed']}/{highest['_total']})")
    callout_top = Inches(table_top + table_h + 0.2)
    if callout_top < Inches(6.5):
        _add_callouts(slide, callouts[:5], callout_top)
    return slide


# ── Color rule helpers (match dashboard thresholds) ──

def _kds_colors(df_raw, df_display):
    """Return cell_colors + bold_cells for KDS table. df_raw has raw numeric values."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    for i, (_, row) in enumerate(df_raw.iterrows()):
        # SOS > 10 → red
        if "SOS" in col_names and pd.notna(row.get("SOS")) and row["SOS"] > 10:
            j = col_names.index("SOS")
            colors[(i, j)] = RED; bolds[(i, j)] = True
        # Adoption < 85 → red
        if "Adoption %" in col_names and pd.notna(row.get("Adoption %")) and row["Adoption %"] < 85:
            j = col_names.index("Adoption %")
            colors[(i, j)] = RED; bolds[(i, j)] = True
        # Make Ahead > 10 → red
        if "Make Ahead %" in col_names and pd.notna(row.get("Make Ahead %")) and row["Make Ahead %"] > 10:
            j = col_names.index("Make Ahead %")
            colors[(i, j)] = RED; bolds[(i, j)] = True
        # Waste > 5 → red
        if "Waste %" in col_names and pd.notna(row.get("Waste %")) and row["Waste %"] > 5:
            j = col_names.index("Waste %")
            colors[(i, j)] = RED; bolds[(i, j)] = True
        # Pre-Bump > 0.5 → red
        if "Pre-Bump %" in col_names and pd.notna(row.get("Pre-Bump %")) and row["Pre-Bump %"] > 0.5:
            j = col_names.index("Pre-Bump %")
            colors[(i, j)] = RED; bolds[(i, j)] = True
        # Adherence < 85 → red
        if "Adherence %" in col_names and pd.notna(row.get("Adherence %")) and row["Adherence %"] < 85:
            j = col_names.index("Adherence %")
            colors[(i, j)] = RED; bolds[(i, j)] = True
    return colors, bolds


def _labor_colors(df_raw, df_display, var_col_name="Labor Var %", var_field="hours_var_pct"):
    """Color the labor-variance column by the bonus threshold (pass <= 1.5%)."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    if var_col_name in col_names:
        j_var = col_names.index(var_col_name)
        for i, (_, row) in enumerate(df_raw.iterrows()):
            var = row.get(var_field)
            if pd.notna(var):
                # Over guide is bad; <= 1.5% passes (bonus target), <= 3% tolerance, else red
                if var <= 1.5:
                    colors[(i, j_var)] = GREEN; bolds[(i, j_var)] = True
                elif var <= 3:
                    colors[(i, j_var)] = ORANGE; bolds[(i, j_var)] = True
                else:
                    colors[(i, j_var)] = RED; bolds[(i, j_var)] = True
    for i, (_, row) in enumerate(df_raw.iterrows()):
        # Labor % > 20% critical red, > 18% orange
        labor_pct = row.get("labor_pct")
        if pd.notna(labor_pct) and "Labor %" in col_names:
            j_lp = col_names.index("Labor %")
            if labor_pct > 20:
                colors[(i, j_lp)] = RED; bolds[(i, j_lp)] = True
            elif labor_pct > 18:
                colors[(i, j_lp)] = ORANGE; bolds[(i, j_lp)] = True
    return colors, bolds


def _smg_colors(df_raw, df_display):
    """Return cell_colors for SMG table."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    for i, (_, row) in enumerate(df_raw.iterrows()):
        if pd.notna(row.get("Dissatisfaction %")) and row["Dissatisfaction %"] > 3:
            if "Dissat %" in col_names:
                j = col_names.index("Dissat %")
                colors[(i, j)] = RED; bolds[(i, j)] = True
        if pd.notna(row.get("Inaccurate Order %")) and row["Inaccurate Order %"] > 3:
            if "Inaccurate %" in col_names:
                j = col_names.index("Inaccurate %")
                colors[(i, j)] = RED; bolds[(i, j)] = True
        if pd.notna(row.get("Greeted with Smile %")) and row["Greeted with Smile %"] < 95:
            if "Greeted %" in col_names:
                j = col_names.index("Greeted %")
                colors[(i, j)] = RED; bolds[(i, j)] = True
    return colors, bolds


def _sales_colors(df_raw, df_display):
    """Return cell_colors for Sales table — Cash O/S off guide if >±$2."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    if "Cash O/S" not in col_names:
        return colors, bolds
    j = col_names.index("Cash O/S")
    for i, (_, row) in enumerate(df_raw.iterrows()):
        cos = row.get("Cash Over/Short")
        if pd.notna(cos) and abs(cos) > 2:
            colors[(i, j)] = RED; bolds[(i, j)] = True
    return colors, bolds


def _flavorlab_colors(df_raw, df_display):
    """Return cell_colors for FlavorLab table."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    if "Completion %" not in col_names:
        return colors, bolds
    j = col_names.index("Completion %")
    for i, (_, row) in enumerate(df_raw.iterrows()):
        val = row.get("Avg_Completion_Rate")
        if pd.notna(val) and val < 95:
            colors[(i, j)] = RED; bolds[(i, j)] = True
    return colors, bolds


def _cogs_colors(df_raw, df_display):
    """Return cell_colors for COGS table."""
    colors, bolds = {}, {}
    col_names = list(df_display.columns)
    if "COGS Variance %" not in col_names:
        return colors, bolds
    j = col_names.index("COGS Variance %")
    for i, (_, row) in enumerate(df_raw.iterrows()):
        val = row.get("cogs_var")
        if pd.notna(val) and abs(val) > 1:
            colors[(i, j)] = RED; bolds[(i, j)] = True
        elif pd.notna(val) and val < 1:
            colors[(i, j)] = GREEN; bolds[(i, j)] = True
    return colors, bolds


def generate_district_ppt(district, store_to_district, districts_config):
    """Generate a complete PPT for the given district. Returns bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    d_stores = _get_district_stores(district, store_to_district)

    # ════════════════════════════════════════
    # SLIDE 1: TITLE / SUMMARY
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    # Full green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GREEN
    bg.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "FL Wingmen"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = f"{district} — Weekly Operations Report"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
    p2.alignment = PP_ALIGN.LEFT

    # Summary stats
    summary_items = []

    # KDS adherence
    kds_path = DATA_DIR / "kds_dinner.csv"
    if kds_path.exists():
        kds = pd.read_csv(kds_path)
        kds["Store No"] = kds["Store No"].astype(str)
        latest_kds = sorted(kds["Period"].unique(), key=lambda x: (int(x[1]), int(x[3])))[-1]
        kds_d = kds[(kds["Period"] == latest_kds) & (kds["Store No"].isin(d_stores))]
        if not kds_d.empty:
            avg_adh = kds_d["Adherence %"].mean()
            summary_items.append(f"KDS Adherence ({latest_kds}): {avg_adh:.0f}%")

    # FlavorLab
    fl_path = DATA_DIR / "flavorlab.csv"
    if fl_path.exists():
        fl = pd.read_csv(fl_path)
        fl["Store No"] = fl["Store No"].astype(str)
        fl_d = fl[fl["Store No"].isin(d_stores)]
        if not fl_d.empty:
            fl_avg = fl_d["Avg_Completion_Rate"].mean()
            summary_items.append(f"FlavorLab Completion: {fl_avg:.0f}%")

    # SMG (most recent quarter available)
    _smg_q = max([q for q in (1, 2, 3, 4) if (DATA_DIR / f"smg_q{q}.csv").exists()], default=2)
    smg_path = DATA_DIR / f"smg_q{_smg_q}.csv"
    smg_qlabel = f"Q{_smg_q}"
    if smg_path.exists():
        smg = pd.read_csv(smg_path)
        smg["Store No"] = smg["Store No"].astype(str)
        smg_d = smg[smg["Store No"].isin(d_stores)]
        if not smg_d.empty and smg_d["Survey Count"].sum() > 0:
            w_dissat = (smg_d["Dissatisfaction %"] * smg_d["Survey Count"]).sum() / smg_d["Survey Count"].sum()
            summary_items.append(f"SMG Dissatisfaction ({smg_qlabel}): {w_dissat:.2f}%")

    summary_items.append(f"Stores: {len(d_stores)}")

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(3))
    tf2 = txBox2.text_frame
    for item in summary_items:
        p = tf2.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    # ════════════════════════════════════════
    # SLIDE 2: SALES PERFORMANCE
    # ════════════════════════════════════════
    sales_path = DATA_DIR / "sales_journal.csv"
    if sales_path.exists():
        sales = pd.read_csv(sales_path)
        sales["Store No"] = sales["Store No"].astype(str)
        all_sales_periods = sorted(sales["Period"].unique(), key=lambda x: (int(x[1]), int(x[3])))
        latest_sales = all_sales_periods[-1]
        prev_sales = all_sales_periods[-2] if len(all_sales_periods) >= 2 else None
        sales_d = sales[(sales["Period"] == latest_sales) & (sales["Store No"].isin(d_stores))].copy()

        # Prior week for WoW delta
        prev_gross_map = {}
        if prev_sales:
            sales_prev = sales[(sales["Period"] == prev_sales) & (sales["Store No"].isin(d_stores))]
            prev_gross_map = dict(zip(sales_prev["Store No"], sales_prev["Gross Sales"]))

        if not sales_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, f"Sales Performance — {latest_sales}", f"{district} · {len(sales_d)} stores")

            sales_d["Online %"] = (sales_d["Online Sales"] / sales_d["Gross Sales"] * 100).round(1)
            raw_tbl = sales_d[["Store No", "Gross Sales", "Net Sales", "Online Sales", "Online %",
                               "Check Avg", "Cash Over/Short"]].copy()
            raw_tbl = raw_tbl.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)

            # WoW gross sales delta
            raw_tbl["Δ Sales"] = raw_tbl.apply(
                lambda r: r["Gross Sales"] - prev_gross_map.get(r["Store No"], r["Gross Sales"])
                if pd.notna(r["Gross Sales"]) else None, axis=1)

            tbl = raw_tbl.copy()
            tbl["Gross Sales"] = tbl["Gross Sales"].apply(lambda x: f"${x:,.0f}" if pd.notna(x) else "—")
            tbl["Net Sales"] = tbl["Net Sales"].apply(lambda x: f"${x:,.0f}" if pd.notna(x) else "—")
            tbl["Online Sales"] = tbl["Online Sales"].apply(lambda x: f"${x:,.0f}" if pd.notna(x) else "—")
            tbl["Online %"] = tbl["Online %"].apply(lambda x: f"{x:.1f}%" if pd.notna(x) else "—")
            tbl["Check Avg"] = tbl["Check Avg"].apply(lambda x: f"${x:,.2f}" if pd.notna(x) else "—")
            tbl["Cash Over/Short"] = tbl["Cash Over/Short"].apply(lambda x: f"${x:+,.2f}" if pd.notna(x) else "—")
            tbl["Δ Sales"] = raw_tbl["Δ Sales"].apply(
                lambda x: f"{'↑' if x > 0 else '↓'} ${abs(x):,.0f}" if pd.notna(x) and x != 0 else ("→ $0" if pd.notna(x) else "—"))
            _sales_prev_label = f"vs {prev_sales}" if prev_sales else "vs Last Wk"
            tbl.columns = ["Store #", "Gross Sales", "Net Sales", "Online $", "Online %",
                           "Check Avg", "Cash O/S", f"Sales ({_sales_prev_label})"]

            cc, bc = _sales_colors(raw_tbl, tbl)
            # Color WoW sales delta: increase=green, decrease=red
            col_names = list(tbl.columns)
            j_ds = len(col_names) - 1  # last column
            for i, (_, row) in enumerate(raw_tbl.iterrows()):
                d = row.get("Δ Sales")
                if pd.notna(d) and d != 0:
                    cc[(i, j_ds)] = GREEN if d > 0 else RED
                    bc[(i, j_ds)] = True

            _add_table(slide, tbl, Inches(0.5), Inches(1.4), SLIDE_WIDTH - Inches(1),
                       Inches(0.3 * (len(tbl) + 1)), cell_colors=cc, bold_cells=bc)

            # Callouts
            callouts = []
            total_gross = sales_d["Gross Sales"].sum()
            prev_total = sum(prev_gross_map.get(s, 0) for s in raw_tbl["Store No"]) if prev_gross_map else 0
            callouts.append(f"Total gross sales: ${total_gross:,.0f}")
            if prev_total > 0:
                delta_total = total_gross - prev_total
                callouts.append(f"WoW change: ${delta_total:+,.0f} vs {prev_sales}")
            cos_bad = sales_d[sales_d["Cash Over/Short"].abs() > 2]
            for _, r in cos_bad.iterrows():
                callouts.append(f"Store {r['Store No']} cash O/S: ${r['Cash Over/Short']:+,.2f} (off guide ±$2)")
            top_table_h = Inches(1.4 + 0.3 * (len(tbl) + 1) + 0.2)
            _add_callouts(slide, callouts[:5], top_table_h)

    # ════════════════════════════════════════
    # SLIDE 3: KDS DASHBOARD
    # ════════════════════════════════════════
    if kds_path.exists():
        kds = pd.read_csv(kds_path)
        kds["Store No"] = kds["Store No"].astype(str)
        all_periods = sorted(kds["Period"].unique(), key=lambda x: (int(x[1]), int(x[3])))
        latest_kds = all_periods[-1]
        prev_kds = all_periods[-2] if len(all_periods) >= 2 else None
        kds_d = kds[(kds["Period"] == latest_kds) & (kds["Store No"].isin(d_stores))].copy()

        # Prior week for WoW deltas
        prev_map_sos, prev_map_adh = {}, {}
        if prev_kds:
            kds_prev = kds[(kds["Period"] == prev_kds) & (kds["Store No"].isin(d_stores))]
            prev_map_sos = dict(zip(kds_prev["Store No"], kds_prev["SOS"]))
            prev_map_adh = dict(zip(kds_prev["Store No"], kds_prev["Adherence %"]))

        if not kds_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, f"KDS Dashboard — {latest_kds}", f"{district} · Fri & Sat Dinner Only")

            raw_kds = kds_d[["Store No", "Store Name", "SOS", "Adoption %", "Make Ahead %",
                             "Waste %", "Pre-Bump %", "Adherence %"]].copy()
            raw_kds = raw_kds.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)

            # Compute WoW deltas
            raw_kds["Δ SOS"] = raw_kds.apply(
                lambda r: r["SOS"] - prev_map_sos.get(r["Store No"], r["SOS"]) if pd.notna(r["SOS"]) else None, axis=1)
            raw_kds["Δ Adh"] = raw_kds.apply(
                lambda r: r["Adherence %"] - prev_map_adh.get(r["Store No"], r["Adherence %"]) if pd.notna(r["Adherence %"]) else None, axis=1)

            tbl = raw_kds.copy()
            tbl["SOS"] = tbl["SOS"].apply(_fmt_sos)
            tbl["Adoption %"] = tbl["Adoption %"].apply(lambda x: f"{x:.1f}" if pd.notna(x) else "—")
            tbl["Make Ahead %"] = tbl["Make Ahead %"].apply(lambda x: f"{x:.1f}" if pd.notna(x) else "—")
            tbl["Waste %"] = tbl["Waste %"].apply(lambda x: f"{x:.2f}" if pd.notna(x) else "—")
            tbl["Pre-Bump %"] = tbl["Pre-Bump %"].apply(lambda x: f"{x:.2f}" if pd.notna(x) else "—")
            tbl["Adherence %"] = tbl["Adherence %"].apply(lambda x: f"{x:.0f}%" if pd.notna(x) else "—")
            # Format deltas with arrows: SOS lower = better (↓ green), Adherence higher = better (↑ green)
            # "→ 0" for no change, "—" only for truly missing data
            def _fmt_delta_sos(x):
                if pd.isna(x): return "—"
                if x == 0: return "→ 0"
                arrow = "↓" if x < 0 else "↑"
                return f"{arrow} {abs(x):.1f}"

            def _fmt_delta_adh(x):
                if pd.isna(x): return "—"
                if x == 0: return "→ 0%"
                arrow = "↑" if x > 0 else "↓"
                return f"{arrow} {abs(x):.0f}%"

            tbl["Δ SOS"] = raw_kds["Δ SOS"].apply(_fmt_delta_sos)
            tbl["Δ Adh"] = raw_kds["Δ Adh"].apply(_fmt_delta_adh)

            _prev_label = f"vs {prev_kds}" if prev_kds else "vs Last Wk"
            tbl.columns = ["Store #", "Store Name", "SOS", "Adoption %", "Make Ahead %",
                           "Waste %", "Pre-Bump %", "Adherence %",
                           f"SOS ({_prev_label})", f"Adh ({_prev_label})"]

            cc, bc = _kds_colors(raw_kds, tbl)
            # Color WoW delta columns: SOS decrease=green increase=red, Adherence increase=green decrease=red
            col_names = list(tbl.columns)
            j_dsos = len(col_names) - 2  # second-to-last column
            j_dadh = len(col_names) - 1  # last column
            for i, (_, row) in enumerate(raw_kds.iterrows()):
                ds = row.get("Δ SOS")
                if pd.notna(ds) and ds != 0:
                    cc[(i, j_dsos)] = GREEN if ds < 0 else RED
                    bc[(i, j_dsos)] = True
                da = row.get("Δ Adh")
                if pd.notna(da) and da != 0:
                    cc[(i, j_dadh)] = GREEN if da > 0 else RED
                    bc[(i, j_dadh)] = True

            # Rows wrap to ~2 lines (long store names + 10 columns), so allocate taller rows
            kds_h = 0.55 + 0.42 * len(tbl)
            _add_table(slide, tbl, Inches(0.5), Inches(1.4), SLIDE_WIDTH - Inches(1),
                       Inches(kds_h), cell_colors=cc, bold_cells=bc,
                       col_widths=[0.8, 2.2, 1.0, 1.1, 1.2, 0.9, 1.1, 1.1, 1.2, 1.2])

            callouts = []
            slow = raw_kds[raw_kds["SOS"] > 10]
            for _, r in slow.iterrows():
                callouts.append(f"Store {r['Store No']} SOS: {_fmt_sos(r['SOS'])} (target < 10:00)")
            low_adopt = raw_kds[raw_kds["Adoption %"] < 85]
            for _, r in low_adopt.iterrows():
                callouts.append(f"Store {r['Store No']} adoption: {r['Adoption %']:.1f}% (target ≥ 85%)")
            low_adh = raw_kds[raw_kds["Adherence %"] < 85]
            for _, r in low_adh.iterrows():
                callouts.append(f"Store {r['Store No']} adherence: {r['Adherence %']:.0f}%")
            if not callouts:
                callouts.append("All stores within KDS targets ✓")
            top_h = Inches(1.4 + kds_h + 0.2)
            _add_callouts(slide, callouts[:5], top_h)

    # ════════════════════════════════════════
    # SLIDE 4: SCHEDULE GUIDE
    # ════════════════════════════════════════
    sched_path = DATA_DIR / "schedule_guide.csv"
    if sched_path.exists():
        sched = pd.read_csv(sched_path)
        sched["Store No"] = sched["Store No"].astype(str)
        sched_periods = sorted(sched["Period"].unique(), key=lambda x: (int(x[1]), int(x[3])))
        latest_sched = sched_periods[-1]
        prev_sched = sched_periods[-2] if len(sched_periods) >= 2 else None
        sched_d = sched[(sched["Period"] == latest_sched) & (sched["Store No"].isin(d_stores))].copy()
        if not sched_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            sub = f"{district} · Weekly forecast & staffing" + (f" · Δ Hrs vs {prev_sched}" if prev_sched else "")
            _add_title_bar(slide, f"Schedule Guide — {latest_sched}", sub)

            # Previous week's hours guide per store, for the +/- delta
            prev_hrs = {}
            if prev_sched:
                ph = sched[(sched["Period"] == prev_sched) & (sched["Store No"].isin(d_stores))]
                prev_hrs = dict(zip(ph["Store No"], ph["Hours Guide"]))

            raw_sched = sched_d[["Store No", "Store Name", "Sales Forecast", "Hours Guide"]].copy()
            raw_sched = raw_sched.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)
            raw_sched["dHrs"] = raw_sched.apply(
                lambda r: (r["Hours Guide"] - prev_hrs[r["Store No"]]) if r["Store No"] in prev_hrs and pd.notna(prev_hrs[r["Store No"]]) else None, axis=1)

            tbl = raw_sched.copy()
            tbl["Sales Forecast"] = tbl["Sales Forecast"].apply(lambda x: f"${x:,.0f}" if pd.notna(x) else "—")
            tbl["Hours Guide"] = tbl["Hours Guide"].apply(lambda x: f"{x:,.0f}" if pd.notna(x) else "—")
            tbl["dHrs"] = tbl["dHrs"].apply(lambda x: f"{x:+,.0f}" if pd.notna(x) else "—")
            tbl = tbl[["Store No", "Store Name", "Sales Forecast", "Hours Guide", "dHrs"]]
            tbl.columns = ["Store #", "Store Name", "Sales Forecast", "Hours Guide", "Δ Hrs (vs LW)"]

            _add_table(slide, tbl, Inches(0.5), Inches(1.4), Inches(9), Inches(0.3 * (len(tbl) + 1)))

            callouts = []
            total_fc = sched_d["Sales Forecast"].sum()
            total_hrs = sched_d["Hours Guide"].sum()
            callouts.append(f"Total forecast: ${total_fc:,.0f} | Total hours: {total_hrs:,.0f}")
            top_store = sched_d.loc[sched_d["Sales Forecast"].idxmax()]
            callouts.append(f"Highest forecast: Store {top_store['Store No']} — ${top_store['Sales Forecast']:,.0f}")
            top_h = Inches(1.4 + 0.3 * (len(tbl) + 1) + 0.2)
            _add_callouts(slide, callouts, top_h)

    # ════════════════════════════════════════
    # SLIDE 5: INTERNAL QSC EVALS
    # ════════════════════════════════════════
    qsc_path = DATA_DIR / "qsc_evals.csv"
    if qsc_path.exists():
        qsc = pd.read_csv(qsc_path)
        qsc["Store No"] = qsc["Store No"].astype(str)
        latest_qsc = sorted(qsc["Period"].unique(), key=lambda x: (int(x[1]), int(x[3])))[-1]
        qsc_d = qsc[(qsc["Period"] == latest_qsc) & (qsc["Store No"].isin(d_stores))].copy()
        if not qsc_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, f"Internal QSC Evals — {latest_qsc}", f"{district} · {len(qsc_d)} evaluations")

            raw_qsc = qsc_d[["Store No", "City", "Date", "Score", "Stars", "Rating"]].copy()
            raw_qsc = raw_qsc.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)

            tbl = raw_qsc.copy()
            tbl.columns = ["Store #", "City", "Date", "Findings", "Stars", "Rating"]

            # Color: < 4 stars → red
            colors, bolds = {}, {}
            col_names = list(tbl.columns)
            j_stars = col_names.index("Stars")
            for i, (_, row) in enumerate(raw_qsc.iterrows()):
                if pd.notna(row.get("Stars")) and row["Stars"] < 4:
                    colors[(i, j_stars)] = RED; bolds[(i, j_stars)] = True

            # Format stars display after computing colors
            tbl["Stars"] = tbl["Stars"].apply(lambda x: f"{'★' * int(x)}" if pd.notna(x) and x > 0 else "—")

            _add_table(slide, tbl, Inches(0.5), Inches(1.4), SLIDE_WIDTH - Inches(1),
                       Inches(0.3 * (len(tbl) + 1)), cell_colors=colors, bold_cells=bolds)

            callouts = []
            five_star = qsc_d[qsc_d["Stars"] == 5]
            callouts.append(f"{len(five_star)}/{len(qsc_d)} evaluations achieved 5 stars")
            below_4 = qsc_d[qsc_d["Stars"] < 4]
            for _, r in below_4.iterrows():
                callouts.append(f"Store {r['Store No']}: {int(r['Stars'])} stars — {r['Rating']}")
            if below_4.empty:
                callouts.append("No stores below 4 stars ✓")
            top_h = Inches(1.4 + 0.3 * (len(tbl) + 1) + 0.2)
            _add_callouts(slide, callouts[:5], top_h)

    # ════════════════════════════════════════
    # SLIDE 6: LABOR DASHBOARD — Latest Week + QTD
    # ════════════════════════════════════════
    forecast_path = DATA_DIR / "forecast.xlsm"
    qtr_ps = [4, 5, 6]   # default Q2
    current_qtr = 2
    if forecast_path.exists():
        raw_labor = pd.read_excel(forecast_path, sheet_name="Forecast_Data")
        r26 = raw_labor[raw_labor["year"] == 2026]
        # Current quarter detection
        latest_p = r26["period"].max()
        qtr_periods = {1: [1,2,3], 2: [4,5,6], 3: [7,8,9], 4: [10,11,12]}
        for q, ps in qtr_periods.items():
            if latest_p in ps:
                current_qtr = q
                qtr_ps = ps
                break

        # Parse store number
        r26 = r26.copy()
        r26["store_num"] = r26["store"].apply(
            lambda x: str(int(str(x).split("-")[0].strip().lstrip("0"))) if "-" in str(x) else str(x))
        r26_d = r26[r26["store_num"].isin(d_stores)]

        if not r26_d.empty:
            # ── Find latest week ──
            latest_period = r26_d["period"].max()
            latest_week_candidates = r26_d[r26_d["period"] == latest_period]
            latest_week = latest_week_candidates["week_d"].max() if "week_d" in r26_d.columns else None

            # Latest week data
            if latest_week is not None:
                lw_data = r26_d[(r26_d["period"] == latest_period) & (r26_d["week_d"] == latest_week)]
            else:
                lw_data = r26_d[r26_d["period"] == latest_period]

            # QTD data
            qtd_data = r26_d[r26_d["period"].isin(qtr_ps)]

            # Latest week period label — week_d format is "2026P06W2", extract "P6W2"
            if latest_week is not None:
                m = re.search(r'P0*(\d+)W(\d+)', str(latest_week))
                lw_label = f"P{m.group(1)}W{m.group(2)}" if m else str(latest_week)
            else:
                lw_label = f"P{latest_period}"

            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, f"Labor Dashboard — {lw_label} & Q{current_qtr} QTD",
                           f"{district} · Variance drives bonus incentive")

            # ── Prior week data for WoW delta ──
            all_weeks = sorted(r26_d["week_d"].unique()) if "week_d" in r26_d.columns else []
            prev_var_map = {}
            if latest_week and len(all_weeks) >= 2:
                lw_idx = list(all_weeks).index(latest_week) if latest_week in all_weeks else -1
                if lw_idx >= 1:
                    prev_week = all_weeks[lw_idx - 1]
                    pw_data = r26_d[r26_d["week_d"] == prev_week]
                    pw_agg = pw_data.groupby("store_num").agg(
                        guide_hours=("guide_hours", "sum"), actual_hours=("actual_crew_hours", "sum"),
                    ).reset_index()
                    # Labor Variance % = column J: actual crew hours vs guide hours
                    pw_agg["variance"] = (pw_agg["actual_hours"] - pw_agg["guide_hours"]) / pw_agg["guide_hours"] * 100
                    prev_var_map = dict(zip(pw_agg["store_num"], pw_agg["variance"]))

            # ── Section 1: Latest Week ──
            _add_section_label(slide, f"Latest Week ({lw_label})", Inches(0.5), Inches(1.2), Inches(5))

            lw_agg = lw_data.groupby("store_num").agg(
                forecast_sales=("forecast_sales", "sum"),
                actual_sales=("actual_sales", "sum"),
                actual_labor=("actual_labor", "sum"),
                schedule_labor=("schedule_labor", "sum"),
                guide_hours=("guide_hours", "sum"),
                actual_hours=("actual_crew_hours", "sum"),
                ovt_hours=("ovt_hours", "sum"),
            ).reset_index()
            lw_agg["labor_pct"] = (lw_agg["actual_labor"] / lw_agg["actual_sales"] * 100)
            lw_agg["sched_pct"] = (lw_agg["schedule_labor"] / lw_agg["forecast_sales"] * 100)
            lw_agg["variance"] = lw_agg["labor_pct"] - lw_agg["sched_pct"]
            # Column J: hours variance vs guide, as %
            lw_agg["hours_var_pct"] = (lw_agg["actual_hours"] - lw_agg["guide_hours"]) / lw_agg["guide_hours"] * 100
            # WoW delta on Labor Variance % (col J: hours vs guide)
            lw_agg["Δ Var"] = lw_agg.apply(
                lambda r: r["hours_var_pct"] - prev_var_map.get(r["store_num"], r["hours_var_pct"]), axis=1)
            lw_agg = lw_agg.sort_values("store_num", key=lambda x: x.astype(int)).reset_index(drop=True)

            lw_tbl = lw_agg[["store_num", "actual_sales", "labor_pct", "sched_pct", "hours_var_pct",
                              "Δ Var", "actual_hours", "variance", "ovt_hours"]].copy()
            lw_display = lw_tbl.copy()
            lw_display["actual_sales"] = lw_display["actual_sales"].apply(lambda x: f"${x:,.0f}")
            lw_display["labor_pct"] = lw_display["labor_pct"].apply(lambda x: f"{x:.1f}%")
            lw_display["sched_pct"] = lw_display["sched_pct"].apply(lambda x: f"{x:.1f}%")
            lw_display["hours_var_pct"] = lw_display["hours_var_pct"].apply(lambda x: f"{x:+.1f}%" if pd.notna(x) else "—")
            lw_display["Δ Var"] = lw_agg["Δ Var"].apply(
                lambda x: f"{'↓' if x < 0 else '↑'} {abs(x):.1f}%" if pd.notna(x) and abs(x) > 0.05 else ("→ 0%" if pd.notna(x) else "—"))
            lw_display["actual_hours"] = lw_display["actual_hours"].apply(lambda x: f"{x:,.0f}")
            lw_display["variance"] = lw_display["variance"].apply(lambda x: f"{x:+.1f}%")
            lw_display["ovt_hours"] = lw_display["ovt_hours"].apply(lambda x: f"{x:,.1f}")
            lw_display.columns = ["Store #", "Actual Sales", "Labor %", "Sched %", "Labor Var %",
                                   "Var (vs Last Wk)", "Hours", "Act vs Sched %", "OT Hours"]

            cc_lw, bc_lw = _labor_colors(lw_agg, lw_display)
            # Color WoW variance delta: decrease=green (improving), increase=red (worsening)
            lw_col_names = list(lw_display.columns)
            j_dvar = lw_col_names.index("Var (vs Last Wk)")
            for i, (_, row) in enumerate(lw_agg.iterrows()):
                dv = row.get("Δ Var")
                if pd.notna(dv) and abs(dv) > 0.05:
                    cc_lw[(i, j_dvar)] = GREEN if dv < 0 else RED
                    bc_lw[(i, j_dvar)] = True

            lw_table_h = max(Inches(0.25 * (len(lw_display) + 1)), Inches(1.5))
            _add_table(slide, lw_display, Inches(0.5), Inches(1.55), Inches(10), lw_table_h,
                       cell_colors=cc_lw, bold_cells=bc_lw)

            # ── Section 2: QTD ──
            qtd_top = Inches(1.55) + lw_table_h + Inches(0.15)
            _add_section_label(slide, f"Quarter-to-Date (Q{current_qtr})", Inches(0.5), qtd_top, Inches(5))

            qtd_agg = qtd_data.groupby("store_num").agg(
                forecast_sales=("forecast_sales", "sum"),
                actual_sales=("actual_sales", "sum"),
                actual_labor=("actual_labor", "sum"),
                schedule_labor=("schedule_labor", "sum"),
                guide_hours=("guide_hours", "sum"),
                actual_hours=("actual_crew_hours", "sum"),
                ovt_hours=("ovt_hours", "sum"),
            ).reset_index()
            qtd_agg["labor_pct"] = (qtd_agg["actual_labor"] / qtd_agg["actual_sales"] * 100)
            qtd_agg["sched_pct"] = (qtd_agg["schedule_labor"] / qtd_agg["forecast_sales"] * 100)
            qtd_agg["variance"] = qtd_agg["labor_pct"] - qtd_agg["sched_pct"]
            # Column J: hours variance vs guide, as %
            qtd_agg["hours_var_pct"] = (qtd_agg["actual_hours"] - qtd_agg["guide_hours"]) / qtd_agg["guide_hours"] * 100
            qtd_agg = qtd_agg.sort_values("store_num", key=lambda x: x.astype(int)).reset_index(drop=True)

            qtd_tbl = qtd_agg[["store_num", "actual_sales", "labor_pct", "sched_pct", "hours_var_pct",
                                "actual_hours", "variance", "ovt_hours"]].copy()
            qtd_display = qtd_tbl.copy()
            qtd_display["actual_sales"] = qtd_display["actual_sales"].apply(lambda x: f"${x:,.0f}")
            qtd_display["labor_pct"] = qtd_display["labor_pct"].apply(lambda x: f"{x:.1f}%")
            qtd_display["sched_pct"] = qtd_display["sched_pct"].apply(lambda x: f"{x:.1f}%")
            qtd_display["hours_var_pct"] = qtd_display["hours_var_pct"].apply(lambda x: f"{x:+.1f}%" if pd.notna(x) else "—")
            qtd_display["actual_hours"] = qtd_display["actual_hours"].apply(lambda x: f"{x:,.0f}")
            qtd_display["variance"] = qtd_display["variance"].apply(lambda x: f"{x:+.1f}%")
            qtd_display["ovt_hours"] = qtd_display["ovt_hours"].apply(lambda x: f"{x:,.1f}")
            qtd_display.columns = ["Store #", "Actual Sales", "Labor %", "Sched %", "Labor Var %", "Hours", "Act vs Sched %", "OT Hours"]

            cc_qtd, bc_qtd = _labor_colors(qtd_agg, qtd_display)
            qtd_table_h = max(Inches(0.25 * (len(qtd_display) + 1)), Inches(1.5))
            _add_table(slide, qtd_display, Inches(0.5), qtd_top + Inches(0.35), Inches(10), qtd_table_h,
                       cell_colors=cc_qtd, bold_cells=bc_qtd)

            # Callouts at bottom
            callout_top = qtd_top + Inches(0.35) + qtd_table_h + Inches(0.1)
            callouts = []
            avg_var_qtd = qtd_agg["hours_var_pct"].mean()
            callouts.append(f"District QTD avg Labor Variance: {avg_var_qtd:+.1f}% (bonus target ≤ 1.5%, hrs vs guide)")
            high_var = qtd_agg[qtd_agg["hours_var_pct"] > 1.5]
            for _, r in high_var.iterrows():
                callouts.append(f"Store {r['store_num']} QTD Labor Variance: {r['hours_var_pct']:+.1f}%")
            total_ot = qtd_agg["ovt_hours"].sum()
            if total_ot > 0:
                callouts.append(f"Total QTD overtime: {total_ot:,.1f} hours")
            if callout_top < Inches(6.8):
                _add_callouts(slide, callouts[:4], callout_top)

    # ════════════════════════════════════════
    # SLIDE 7: COGS VARIANCE (QTD)
    # ════════════════════════════════════════
    cogs_path = DATA_DIR / "cogs_variance.csv"
    if cogs_path.exists():
        cogs = pd.read_csv(cogs_path)
        cogs["Store No"] = cogs["Store No"].astype(str)
        cogs["_p"] = cogs["Period"].str.extract(r"P(\d+)").astype(int)
        cogs_qtr = cogs[cogs["_p"].isin(qtr_ps)]
        cogs_d = cogs_qtr[cogs_qtr["Store No"].isin(d_stores)]
        if not cogs_d.empty:
            cogs_avg = cogs_d.groupby(["Store No", "Store Name"]).agg(
                nbo_actual=("NBO Actual %", "mean"),
                food_theo=("Food Theoretical %", "mean"),
                cogs_actual=("COGS Actual %", "mean"),
                cogs_var=("COGS Variance %", "mean"),
            ).reset_index()
            cogs_avg = cogs_avg.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)

            slide = prs.slides.add_slide(blank_layout)
            cogs_periods = sorted(cogs_qtr["Period"].unique())
            _add_title_bar(slide, f"COGS Variance — {', '.join(cogs_periods)}", f"{district} · QTD")

            tbl = cogs_avg[["Store No", "Store Name", "food_theo", "cogs_actual", "cogs_var"]].copy()
            tbl_display = tbl.copy()
            tbl_display["food_theo"] = tbl_display["food_theo"].apply(lambda x: f"{x:.2f}%")
            tbl_display["cogs_actual"] = tbl_display["cogs_actual"].apply(lambda x: f"{x:.2f}%")
            tbl_display["cogs_var"] = tbl_display["cogs_var"].apply(lambda x: f"{x:+.2f}%")
            tbl_display.columns = ["Store #", "Store Name", "Food Theo %", "COGS Actual %", "COGS Variance %"]

            cc, bc = _cogs_colors(cogs_avg, tbl_display)
            _add_table(slide, tbl_display, Inches(0.5), Inches(1.4), Inches(9),
                       Inches(0.3 * (len(tbl_display) + 1)), cell_colors=cc, bold_cells=bc)

            callouts = []
            avg_var = cogs_avg["cogs_var"].mean()
            callouts.append(f"District avg COGS variance: {avg_var:+.2f}%")
            high_cogs = cogs_avg[cogs_avg["cogs_var"] > 1]
            for _, r in high_cogs.iterrows():
                callouts.append(f"Store {r['Store No']}: COGS variance {r['cogs_var']:+.2f}% (target ≤ 1%)")
            if high_cogs.empty:
                callouts.append("All stores within COGS target ✓")
            top_h = Inches(1.4 + 0.3 * (len(tbl_display) + 1) + 0.2)
            _add_callouts(slide, callouts[:5], top_h)

    # ════════════════════════════════════════
    # SLIDE 8: SMG (QTD)
    # ════════════════════════════════════════
    if smg_path.exists():
        smg = pd.read_csv(smg_path)
        smg["Store No"] = smg["Store No"].astype(str)
        smg_d = smg[smg["Store No"].isin(d_stores)].copy()
        smg_d = smg_d.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)
        if not smg_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, f"SMG — Guest Satisfaction ({smg_qlabel} QTD)",
                           f"{district} · {smg_d['Survey Count'].sum():,} surveys")

            raw_smg = smg_d[["Store No", "Store Name", "Survey Count", "Dissatisfaction %",
                             "Inaccurate Order %", "Greeted with Smile %"]].copy()
            raw_smg = raw_smg.reset_index(drop=True)

            tbl = raw_smg.copy()
            tbl["Dissatisfaction %"] = tbl["Dissatisfaction %"].apply(lambda x: f"{x:.2f}%")
            tbl["Inaccurate Order %"] = tbl["Inaccurate Order %"].apply(lambda x: f"{x:.2f}%")
            tbl["Greeted with Smile %"] = tbl["Greeted with Smile %"].apply(lambda x: f"{x:.1f}%")
            tbl.columns = ["Store #", "Store Name", "Surveys", "Dissat %", "Inaccurate %", "Greeted %"]

            cc, bc = _smg_colors(raw_smg, tbl)
            _add_table(slide, tbl, Inches(0.5), Inches(1.4), Inches(10),
                       Inches(0.3 * (len(tbl) + 1)), cell_colors=cc, bold_cells=bc)

            callouts = []
            w_dissat = (smg_d["Dissatisfaction %"].astype(float) * smg_d["Survey Count"]).sum() / smg_d["Survey Count"].sum()
            callouts.append(f"District weighted dissatisfaction: {w_dissat:.2f}% (target ≤ 3%)")
            high_dissat = smg_d[smg_d["Dissatisfaction %"].astype(float) > 5]
            for _, r in high_dissat.iterrows():
                callouts.append(f"Store {r['Store No']}: {r['Dissatisfaction %']}% dissatisfaction")
            if high_dissat.empty:
                callouts.append("All stores below 5% dissatisfaction ✓")
            top_h = Inches(1.4 + 0.3 * (len(tbl) + 1) + 0.2)
            _add_callouts(slide, callouts[:5], top_h)

    # ════════════════════════════════════════
    # SLIDE 9: FLAVORLAB
    # ════════════════════════════════════════
    if fl_path.exists():
        fl = pd.read_csv(fl_path)
        fl["Store No"] = fl["Store No"].astype(str)
        fl_d = fl[fl["Store No"].isin(d_stores)].copy()
        fl_d = fl_d.sort_values("Store No", key=lambda x: x.astype(int)).reset_index(drop=True)
        if not fl_d.empty:
            slide = prs.slides.add_slide(blank_layout)
            _add_title_bar(slide, "FlavorLab — Training Completion",
                           f"{district} · {fl_d['Employees'].sum()} employees")

            raw_fl = fl_d[["Store No", "Store Name", "Employees", "Total_Courses",
                           "Completions", "Incomplete", "Avg_Completion_Rate"]].copy()
            raw_fl = raw_fl.reset_index(drop=True)

            tbl = raw_fl.copy()
            tbl["Avg_Completion_Rate"] = tbl["Avg_Completion_Rate"].apply(lambda x: f"{x:.1f}%")
            tbl.columns = ["Store #", "Store Name", "Employees", "Courses",
                           "Completed", "Incomplete", "Completion %"]

            cc, bc = _flavorlab_colors(raw_fl, tbl)
            _add_table(slide, tbl, Inches(0.5), Inches(1.4), Inches(10),
                       Inches(0.3 * (len(tbl) + 1)), cell_colors=cc, bold_cells=bc)

            callouts = []
            avg_fl = fl_d["Avg_Completion_Rate"].mean()
            callouts.append(f"District avg completion: {avg_fl:.1f}% (target 100%)")
            low_fl = fl_d[fl_d["Avg_Completion_Rate"] < 95]
            for _, r in low_fl.iterrows():
                callouts.append(f"Store {r['Store No']}: {r['Avg_Completion_Rate']:.1f}% — {int(r['Incomplete'])} incomplete courses")
            if low_fl.empty:
                callouts.append("All stores above 95% completion ✓")
            top_h = Inches(1.4 + 0.3 * (len(tbl) + 1) + 0.2)
            _add_callouts(slide, callouts[:5], top_h)

    # ════════════════════════════════════════
    # SLIDE 10: SCORECARD — QTD Adherence Summary (Concluding Slide)
    # ════════════════════════════════════════
    # Scorecard — one concluding slide per quarter that has a QSC inspection file (Q2 + Q3).
    _sc_qmap = {1: [1, 2, 3], 2: [4, 5, 6], 3: [7, 8, 9], 4: [10, 11, 12]}
    sc_quarters = sorted([q for q in (1, 2, 3, 4) if (DATA_DIR / f"qsc_inspection_q{q}.csv").exists()])
    if not sc_quarters:
        sc_quarters = [current_qtr]

    # Define all checks (same as dashboard scorecard)
    check_defs = [
        ("KDS", "SOS < 10"),
        ("KDS", "Adopt ≥ 85%"),
        ("KDS", "MakeAhd ≤ 10%"),
        ("KDS", "PreBump ≤ 0.5%"),
        ("Labor", "Var ≤ 1.5%"),
        ("SMG", "Dissat ≤ 3%"),
        ("SMG", "Inacc ≤ 3%"),
        ("QSC", "5 Stars"),
        ("FlavorLab", "≥ 95%"),
        ("COGS", "Var ≤ 1%"),
    ]
    check_labels = [f"{cat}: {lbl}" for cat, lbl in check_defs]

    # Quarter-independent data, loaded once
    _sc_kds = None
    if kds_path.exists():
        _sc_kds = pd.read_csv(kds_path)
        _sc_kds["Store No"] = _sc_kds["Store No"].astype(str)
        _sc_kds["_p"] = _sc_kds["Period"].str.extract(r"P(\d+)").astype(int)

    _sc_labor = None
    if forecast_path.exists():
        try:
            _lr = pd.read_excel(forecast_path, sheet_name="Forecast_Data")
            _sc_labor = _lr[_lr["year"] == 2026].copy()
            _sc_labor["store_num"] = _sc_labor["store"].apply(
                lambda x: str(int(str(x).split("-")[0].strip().lstrip("0"))) if "-" in str(x) else str(x))
        except Exception:
            pass

    _sc_fl = None
    if fl_path.exists():
        _sc_fl = pd.read_csv(fl_path)
        _sc_fl["Store No"] = _sc_fl["Store No"].astype(str)

    _sc_cogs = None
    if cogs_path.exists():
        _sc_cogs = pd.read_csv(cogs_path)
        _sc_cogs["Store No"] = _sc_cogs["Store No"].astype(str)
        _sc_cogs["_p"] = _sc_cogs["Period"].str.extract(r"P(\d+)").astype(int)

    # One concluding scorecard slide per quarter (chronological: Q2 then Q3)
    for _scq in sc_quarters:
        _sc_smg = None
        _smg_p = DATA_DIR / f"smg_q{_scq}.csv"
        if _smg_p.exists():
            _sc_smg = pd.read_csv(_smg_p)
            _sc_smg["Store No"] = _sc_smg["Store No"].astype(str)
        _sc_qi = None
        _qi_p = DATA_DIR / f"qsc_inspection_q{_scq}.csv"
        if _qi_p.exists():
            _sc_qi = pd.read_csv(_qi_p)
            _sc_qi["Store No"] = _sc_qi["Store No"].astype(str)
        _add_scorecard_slide(prs, blank_layout, district, d_stores, _scq, _sc_qmap[_scq],
                             check_labels, _sc_kds, _sc_labor, _sc_smg, _sc_qi, _sc_fl, _sc_cogs)

    # ── Save to bytes ──
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
